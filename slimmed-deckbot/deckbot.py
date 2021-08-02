from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches


prs = Presentation()

# Choose layout
layout = prs.slide_layouts[0] 

# First slide
slide = prs.slides.add_slide(layout)
title = slide.shapes.title 
subtitle = slide.placeholders[1]

# Add text
title.text = "This is the title"
subtitle.text = "This is the subtitle"

# Second Slide
slide = prs.slides.add_slide(layout)

chart_data = ChartData()
chart_data.categories = ['Column 1', 'Column 2', 'Column 3']
chart_data.add_series('Hey',    (10.5, 5.5, 17.5))
chart_data.add_series('there',    (25.5, 40.3, 30.7))
chart_data.add_series('reader', (5.2, 10.3, 8.4))

x,y,cx,cy=Inches(1),Inches(2),Inches(5),Inches(7)

# Add Chart
chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS_STACKED,x,y,cx,cy,chart_data).chart

prs.save('test.pptx')